import os
import argparse
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
from slugify import slugify


def add_verse_slide(pres, verse):
    # Title and content layout
    layout = pres.slide_layouts[5]
    slide = pres.slides.add_slide(layout)

    # Add verse text

    # Split string on newlines and strip leading (and trailing) whitespace
    verse_paragraphs = list(map(lambda v: v.strip(), verse.split('\n')))

    verse_text = ''

    # Re-create the verse text for the slide paragraph
    for i, vp in enumerate(verse_paragraphs):
        verse_text += vp

        if i != len(verse_paragraphs) - 1:
            verse_text += '\n'

    slide.shapes.add_textbox(Inches(0.125), Inches(2.5), Inches(13.125), Inches(5))

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text_frame = shape.text_frame

        # Clear out the default empty paragraph
        text_frame.clear()
        text_frame.left = Inches(0)
        text_frame.width = Inches(13)
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.level = 0
        # p.text = verse_text

        run = p.add_run()
        run.text = verse_text
        font = run.font
        font.name = 'Lucida Sans'
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)

    # Add song title as slide title
    slide.shapes.title.text = f'{title} #{song_number}'
    tf = slide.shapes.title.text_frame
    tfp = tf.paragraphs[0]
    tfp.alignment = PP_ALIGN.CENTER
    tfpr = tfp.runs[0]
    tfprf = tfpr.font
    tfprf.name = 'Lucida Sans'
    tfprf.size = Pt(48)
    tfprf.color.rgb = RGBColor(255, 255, 255)
    tfprf.bold = True
    tfprf.italic = True

    # Add Anthem logo to right of slide
    slide.shapes.add_picture('anthem-logo.png', Inches(11.75), Inches(0))


if __name__ == '__main__':
    print('Generating powerpoint presentations...')

    parser = argparse.ArgumentParser()
    parser.add_argument('src_dir')
    parser.add_argument('out_dir')
    parser.add_argument('pptx_base')
    args = parser.parse_args()

    src_dir = args.src_dir
    output_dir = args.out_dir
    base_pptx = args.pptx_base

    if not os.path.exists(src_dir) or not os.path.exists(output_dir) or not os.path.exists(base_pptx):
        if not os.path.exists(src_dir):
            print(f'Provided src_dir of {src_dir} does not exist.')

        if not os.path.exists(output_dir):
            print(f'Provided out_dir of {output_dir} does not exist.')

        if not os.path.exists(base_pptx):
            print(f'Provided pptx_base of {base_pptx} does not exist.')

        exit(-1)

    hymn_xml_files = os.listdir(src_dir)

    for f in hymn_xml_files:
        if not f.endswith('.xml'):
            continue

        print(f'\tCreating pptx for file: {os.path.join(src_dir, f)}')

        with open(os.path.join(src_dir, f)) as fp:
            soup = BeautifulSoup(fp, 'lxml-xml')

            title = soup.find('Title').text.strip()
            song_number = soup.find('SongNumber').text.strip()
            verses = list(map(lambda v: v.text.strip(), soup.findAll('Text')))
            author = soup.find('Author').text.strip()

            # Create a pptx file using the base pptx so we can use its theme
            pres = Presentation(base_pptx)

            # Delete the first slide (yuck). See https://stackoverflow.com/a/70969310
            for i in range(len(pres.slides) - 1, -1, -1):
                rId = pres.slides._sldIdLst[i].rId
                pres.part.drop_rel(rId)
                del pres.slides._sldIdLst[i]

            for verse in verses:
                add_verse_slide(pres, verse)

            # Add the last slide

            # Blank layout
            layout = pres.slide_layouts[6]

            end_slide = pres.slides.add_slide(layout)
            end_slide.shapes.add_textbox(Inches(0.125), Inches(7), Inches(13.125), Inches(0.5))

            for shape in end_slide.shapes:
                if not shape.has_text_frame:
                    continue

                text_frame = shape.text_frame

                # Clear out the default empty paragraph
                text_frame.clear()
                text_frame.left = Inches(0)
                text_frame.width = Inches(13)
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                p.level = 0
                # p.text = verse_text

                run = p.add_run()
                run.text = author
                font = run.font
                font.name = 'Calibri'
                font.size = Pt(14)
                font.color.rgb = RGBColor(255, 255, 255)

            filename = slugify(f'{song_number.zfill(3)} {title}')
            filepath = os.path.join(output_dir, filename)
            pres.save(f'{filepath}.pptx')

    print('Done!')
